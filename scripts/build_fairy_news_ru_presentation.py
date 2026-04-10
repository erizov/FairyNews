#!/usr/bin/env python3
"""Сборка PPTX для защиты по материалам ``docs/fairy_news_booklet_ru``.

Скрипт заполняет стандартный master PowerPoint. При необходимости тему
можно применить вручную после открытия файла в PowerPoint.

Скриншоты берутся из ``docs/fairy_news_booklet/png`` (съёмка Playwright).
Если папка пуста: ``python scripts/build_fairy_news_booklet.py --port 8765``
(сервер, ``FAIRYNEWS_SAVE_REPORTS=1``) или ``… --capture`` у этого скрипта.

Пример::

    python scripts/build_fairy_news_ru_presentation.py
    python scripts/build_fairy_news_ru_presentation.py --capture --port 8765
"""

from __future__ import annotations

import argparse
import subprocess
import sys
from pathlib import Path
from typing import Any

ROOT = Path(__file__).resolve().parent.parent
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

_OUT_DIR = ROOT / "docs" / "fairy_news_booklet_ru"
_DEFAULT_NAME = "Fairy_News_RU_Defense.pptx"
_PNG_DIR = ROOT / "docs" / "fairy_news_booklet" / "png"
# Те же имена, что даёт ``scripts/build_fairy_news_booklet.py``.
_WORKFLOW_SCREENSHOTS: tuple[tuple[str, str], ...] = (
    ("Сценарий: шаг 1 — текст новости", "02_step1_news_ready.png"),
    ("Сценарий: шаг 2 — пресет корпуса (RAG)", "03_step2_preset.png"),
    (
        "Сценарий: шаг 3 — сказка, вопрос, отчёты",
        "04_step3_pipeline_result.png",
    ),
)
_REPORT_SCREENSHOTS: tuple[tuple[str, str], ...] = (
    ("Отчёты: таблица прогонов", "05_reports_index_full.png"),
    ("Отчёты: детальный прогон", "06_detail_after_id_click.png"),
    ("Отчёты: журнал LLM (trace)", "07_llm_log_after_trace_click.png"),
)


def _slide_title_content(
    prs: Any,
    title: str,
    body: str,
    *,
    body_pt: int = 20,
    title_pt: int | None = None,
) -> None:
    from pptx.util import Pt

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    if title_pt is not None:
        for tp in slide.shapes.title.text_frame.paragraphs:
            tp.font.size = Pt(title_pt)
    ph = slide.placeholders[1]
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = body
    p0.font.size = Pt(body_pt)
    p0.level = 0


def _slide_bullets(prs: Any, title: str, bullets: list[str], *, level0_pt: int = 20) -> None:
    from pptx.util import Pt

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for b in bullets:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.text = b
        para.level = 0
        para.font.size = Pt(level0_pt)
    tf.margin_bottom = Pt(12)


def _add_screenshot_slide(prs: Any, title: str, image_path: Path) -> bool:
    """Слайд «только заголовок»: картинка целиком в одной странице, без обрезки."""
    from pptx.util import Inches

    if not image_path.is_file():
        return False
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    max_w_in = 12.35
    max_h_in = 5.72
    top_in = 1.06

    iw_px: int | None = None
    ih_px: int | None = None
    try:
        from PIL import Image

        with Image.open(image_path) as im:
            iw_px, ih_px = im.size
    except OSError:
        iw_px, ih_px = None, None

    if iw_px and ih_px and iw_px > 0 and ih_px > 0:
        ar = ih_px / float(iw_px)
        # Height-first: гарантируем «влезание» по высоте на один слайд.
        h_in = max_h_in
        w_in = h_in / ar
        if w_in > max_w_in:
            w_in = max_w_in
            h_in = w_in * ar
    else:
        w_in = max_w_in
        h_in = max_h_in

    slide_w_in = float(prs.slide_width) / 914400.0
    left_in = max(0.25, (slide_w_in - w_in) / 2.0)
    slide.shapes.add_picture(
        str(image_path),
        Inches(left_in),
        Inches(top_in),
        width=Inches(w_in),
        height=Inches(h_in),
    )
    return True


def _add_workflow_and_report_screenshots(prs: Any) -> tuple[int, list[str]]:
    """Три кадра сценария и три кадра отчётов. Пропускает отсутствующие файлы."""
    added = 0
    missing: list[str] = []
    for title, name in _WORKFLOW_SCREENSHOTS:
        p = _PNG_DIR / name
        if _add_screenshot_slide(prs, title, p):
            added += 1
        else:
            missing.append(name)
    for title, name in _REPORT_SCREENSHOTS:
        p = _PNG_DIR / name
        if _add_screenshot_slide(prs, title, p):
            added += 1
        else:
            missing.append(name)
    return added, missing


def _add_collage_slide(prs: Any) -> None:
    from pptx.util import Inches

    from app.api_schemas import (
        _COLLAGE_TILES,
        _collage_assets_dir,
        _resolved_collage_file,
    )

    pic_ext = frozenset({".jpg", ".jpeg", ".png", ".webp"})
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Коллаж интерфейса Fairy News"

    left = Inches(0.5)
    top = Inches(1.35)
    col_w = Inches(1.35)
    row_h = Inches(1.05)
    gap_x = Inches(0.12)
    gap_y = Inches(0.1)
    base = _collage_assets_dir()

    for idx, (stem, _alt, _m) in enumerate(_COLLAGE_TILES):
        fn, _ = _resolved_collage_file(stem)
        path = base / fn
        row, col = divmod(idx, 4)
        x = left + col * (col_w + gap_x)
        y = top + row * (row_h + gap_y)
        if path.is_file() and path.suffix.lower() in pic_ext:
            slide.shapes.add_picture(str(path), x, y, width=col_w, height=row_h)


def _build_presentation(out_path: Path) -> list[str]:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        print("pip install python-pptx", file=sys.stderr)
        raise SystemExit(1) from None

    prs = Presentation()
    from pptx.util import Inches as _In

    prs.slide_width = int(_In(13.333))
    prs.slide_height = int(_In(7.5))

    opening = (
        "Проблема в том, что у поколения ЕГЭ снижается привычка к бумажной "
        "книге: дети не тренируют изложение — краткий пересказ прочитанного, "
        "выделение главного, отбрасывание побочных сюжетных линий и деталей "
        "антуража, не влияющих на основной посыл автора.\n\n"
        "Боюсь, всё закончится неинтеллигентно, как в сказке той со птицей и "
        "лисицей, далее по тексту…\n\n"
        "Чтобы вернуть детям радость познания и умение обобщать и сравнивать, "
        "а родителям — радость общения с детьми, я и выбрал эту тему.\n\n"
        "Здесь сказки напоминают новости, а новости — сказки; они взаимно "
        "обогащают друг друга.\n\n"
        "• Сопоставление формулировок новостей с устойчивыми сюжетными ходами "
        "фольклора.\n"
        "• Спокойный «книжный» интерфейс: коллаж, переход от хроники к "
        "проверенным формам рассказа."
    )

    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Fairy News"
    for tp in s.shapes.title.text_frame.paragraphs:
        tp.font.size = Pt(50)
    sub = s.placeholders[1]
    sub.text = (
        "Новости через сказку: школа · родительские чаты · Хабр\n"
        "Дипломный проект"
    )
    for p in sub.text_frame.paragraphs:
        p.font.size = Pt(24)

    _slide_title_content(
        prs,
        "Зачем Fairy News",
        opening,
        body_pt=15,
        title_pt=38,
    )

    three_lines = (
        "50+ лет назад фантасты братья Стругацкие писали, что человечество "
        "разделится на две неравные части по неизвестному признаку. Сейчас мы "
        "этот признак видим во всей «славе» его: чтение и понимание длинных "
        "текстов, причём длинных определялось ещё недавно как 7+ страниц, "
        "сейчас уже 5+ страниц и «прогресс» не остановим.\n\n"
        "Не смею утверждать, что мой проект поможет читать длинные тексты, но "
        "заинтересует — это точно. Ведь нет ничего интереснее, чем сказки на "
        "новый лад: иванушка дурачок, кощей бессмертный, василиса премудрая, "
        "скупой богач, глупый мудрец, как интересно тасуется колода…\n\n"
        "• Ребёнку и подростку — узоры сказки и вопрос «что дальше?» без "
        "страха ошибиться.\n"
        "• Взрослому — дистанция и архетипы; факты проверяются вне сказки.\n"
        "• Семье — общий язык: «это как в сказке про…», спор мягче."
    )
    _slide_title_content(
        prs,
        "Неизвестный признак",
        three_lines,
        body_pt=14,
        title_pt=38,
    )

    _slide_bullets(
        prs,
        "Где заходит лучше всего",
        [
            "Школа — 10–20 мин., нейтральная новость, пресет, сюжетное созвучие.",
            "Родительские чаты — короткий пост + демо; один вопрос ребёнку.",
            "Хабр — FastAPI, RAG/Chroma, мультиагентный пайплайн, SQLite, trace.",
        ],
    )

    _slide_bullets(
        prs,
        "Мини-кейс: одна новость — три угла",
        [
            "Новость: мост в посёлке, зима, помощь жителей.",
            "Сказка — мотив общего дела и награды за труд.",
            "Вопрос ребёнку и тема для взрослого — через узнаваемый сюжет.",
        ],
    )

    _slide_bullets(
        prs,
        "Глоссарий (кратко)",
        [
            "Пресет — угол поиска по корпусу сказок.",
            "RAG — релевантные фрагменты в промпт.",
            "Агент — шаг LLM: новость, сказка, аудит, QA.",
        ],
    )

    _slide_bullets(
        prs,
        "Честность и границы",
        [
            "Прогноз — учебная игра шаблонами, не гадание.",
            "Сказка может исказить факты; сверка с источниками и отчёт/trace.",
        ],
    )

    diagram = (
        "Пайплайн (схема):\n"
        "Новость / RSS + пресет → RAG → структура новости → черновик "
        "сказки → аудит + вопрос–ответ → отчёт и журнал LLM."
    )
    _slide_title_content(prs, "Алгоритм", diagram, body_pt=17)

    _slide_bullets(
        prs,
        "Шаги",
        [
            "1. Ввод новости или карточки RSS.",
            "2. Пресет — поиск по корпусу сказок.",
            "3. Агенты: структура → сказка → аудит → вопрос–ответ.",
            "4. RAG в текст.",
            "5. Отчёт и trace по шагам.",
        ],
        level0_pt=19,
    )

    n_shots, miss_shots = _add_workflow_and_report_screenshots(prs)
    if n_shots == 0:
        hint = (
            "Ожидались PNG в каталоге docs/fairy_news_booklet/png (шесть "
            "файлов сценария и отчётов).\n\n"
            "Поднимите сервер с FAIRYNEWS_SAVE_REPORTS=1 и выполните:\n"
            "python scripts/build_fairy_news_booklet.py --port 8765\n\n"
            "либо пересоберите презентацию с ключом --capture."
        )
        _slide_title_content(
            prs,
            "Скриншоты: снимите интерфейс",
            hint,
            body_pt=14,
        )
    elif miss_shots:
        tail = "Отсутствуют: " + ", ".join(miss_shots)
        _slide_title_content(prs, "Пропущенные кадры", tail, body_pt=16)

    _add_collage_slide(prs)

    _slide_bullets(
        prs,
        "Защита: регламент (этап 5)",
        [
            "До 20 мин., очно или дистанционно; свободная подача + слайды.",
            "Цели, задачи, база знаний, сбор и обработка, блок-схема, "
            "алгоритм, эксперименты, трудности и решения, результат, выводы.",
            "Демонстрация: код запущен, сервер и модель доступны.",
        ],
        level0_pt=17,
    )

    _slide_bullets(
        prs,
        "Проблемы и решения",
        [
            "RAG не закрывает весь контекст — локальная подгрузка из .txt, "
            ".pdf и др. в корпус или в промпт.",
            "Связь с LLM — smoke-тест app/llm_connect_try.py (.env, тот же API).",
            "Модель — по результатам тестов.",
            "Аудит в мультиагентной цепочке — целостность и релевантность.",
            "Скрипты start/stop/restart_web (.ps1 / .sh) для демо.",
        ],
        level0_pt=16,
    )

    closing = (
        "Интересного вам чтения, дорогие дети и уважаемые товарищи взрослые.\n\n"
        "Ведь медицинская энциклопедия определяет ум как умение находить "
        "связь между вещами.\n\n"
        "Примечание: та же энциклопедия определяет шизофрению как умение "
        "находить связи там, где их нет — но это уже другая история и не цель "
        "нашего исследования."
    )
    _slide_title_content(prs, "Спасибо за внимание", closing, body_pt=15)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return miss_shots


def main() -> None:
    parser = argparse.ArgumentParser(
        description="PPTX для защиты по материалам RU-буклета.",
    )
    parser.add_argument(
        "-o",
        "--output",
        type=Path,
        default=None,
        help=f"путь к .pptx (по умолчанию: {_OUT_DIR / _DEFAULT_NAME})",
    )
    parser.add_argument(
        "--capture",
        action="store_true",
        help="перед сборкой вызвать build_fairy_news_booklet.py (нужен сервер)",
    )
    parser.add_argument("--host", default="127.0.0.1")
    parser.add_argument("--port", type=int, default=8765)
    args = parser.parse_args()
    out = args.output or (_OUT_DIR / _DEFAULT_NAME)
    if args.capture:
        booklet_script = ROOT / "scripts" / "build_fairy_news_booklet.py"
        subprocess.run(
            [
                sys.executable,
                str(booklet_script),
                "--host",
                args.host,
                "--port",
                str(args.port),
            ],
            check=True,
        )
    missing = _build_presentation(out)
    if missing:
        print(
            "WARN: нет PNG (часть слайдов пропущена):",
            ", ".join(missing),
            file=sys.stderr,
        )
    print("Saved:", out)


if __name__ == "__main__":
    main()
